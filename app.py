import os
import openai
import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates import css, bot_template, user_template

# Load environment variables
load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")
openai.api_base = os.getenv("BASE_API_PATH")
openai.default_headers = {"x-foo": "true"}

def get_docx_text(docx_docs):
    """Extract text from Word documents"""
    from docx import Document
    text = ""
    for doc in docx_docs:
        document = Document(doc)
        text += "\n".join(paragraph.text for paragraph in document.paragraphs)
        text += "\n"
    return text

def get_pptx_text(pptx_docs):
    """Extract text from PowerPoint documents"""
    from pptx import Presentation
    text = ""
    for ppt in pptx_docs:
        presentation = Presentation(ppt)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

def get_excel_text(excel_docs):
    """Extract text from Excel documents"""
    import pandas as pd
    return "\n".join(pd.read_excel(excel).to_string() + "\n" for excel in excel_docs)

def get_txt_text(txt_docs):
    """Extract text from text files"""
    text = ""
    for txt in txt_docs:
        with open(txt, 'r', encoding='utf-8') as file:
            text += file.read() + "\n"
    return text

def get_pdf_text(pdf_docs):
    """Extract text from PDF documents"""
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        text += "".join(page.extract_text() for page in pdf_reader.pages)
    return text

def get_csv_text(csv_docs):
    """Extract text from CSV files"""
    import pandas as pd
    return "\n".join(pd.read_csv(csv).to_string() + "\n" for csv in csv_docs)

def get_file_text(files):
    """Extract text based on file type"""
    text = ""
    file_handlers = {
        '.pdf': get_pdf_text,
        '.docx': get_docx_text,
        '.pptx': get_pptx_text,
        '.xlsx': get_excel_text,
        '.xls': get_excel_text,
        '.txt': get_txt_text,
        '.csv': get_csv_text
    }
    
    for file in files:
        file_extension = os.path.splitext(file.name)[1].lower()
        if file_extension in file_handlers:
            text += file_handlers[file_extension]([file])
        else:
            st.error(f"Unsupported file type: {file_extension}")
    return text

def get_text_chunks(text):
    """Split text into chunks"""
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    return text_splitter.split_text(text)

def get_vectorstore(text_chunks):
    """Create vector store"""
    embeddings = OpenAIEmbeddings(openai_api_key=os.getenv("OPENAI_API_KEY"))
    return FAISS.from_texts(texts=text_chunks, embedding=embeddings)

def get_conversation_chain(vectorstore):
    """Create conversation chain"""
    llm = ChatOpenAI(openai_api_key=os.getenv("OPENAI_API_KEY"))
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True)
    return ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )

def handle_userinput(user_question):
    """Handle user input"""
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']
    
    for i, message in enumerate(st.session_state.chat_history):
        template = user_template if i % 2 == 0 else bot_template
        st.write(template.replace("{{MSG}}", message.content), unsafe_allow_html=True)

def main():
    """Main function"""
    st.set_page_config(page_title="Intelligent Document Q&A System", page_icon=":books:")
    st.write(css, unsafe_allow_html=True)
    
    # Initialize session state
    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None
        
    st.header("Intelligent Document Q&A System :books:")
    
    # Handle user input
    user_question = st.text_input("Enter your question:")
    if user_question:
        handle_userinput(user_question)
        
    # Sidebar file upload
    with st.sidebar:
        st.subheader("Upload Documents")
        uploaded_files = st.file_uploader(
            "Upload PDF, Word, PowerPoint, Excel, TXT, or CSV files",
            type=['pdf', 'docx', 'pptx', 'xlsx', 'xls', 'txt', 'csv'],
            accept_multiple_files=True
        )
        
        if st.button("Process"):
            with st.spinner("Processing documents..."):
                raw_text = get_file_text(uploaded_files)
                text_chunks = get_text_chunks(raw_text)
                vectorstore = get_vectorstore(text_chunks)
                st.session_state.conversation = get_conversation_chain(vectorstore)

if __name__ == '__main__':
    main()
